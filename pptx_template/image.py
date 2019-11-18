#
# coding=utf-8

import logging
import os.path

from six import string_types

import pptx_template.pyel as pyel
import pptx_template.text as txt

from pptx.shapes.picture import Picture
from pptx.shapes.graphfrm import GraphicFrame
from pptx.util import Pt

log = logging.getLogger()


def select_all_table_shapes(slide):
    return [s for s in slide.shapes if isinstance(s, GraphicFrame) and s.shape_type == 19]


def replace_all_img_in_table(shape, model, slide):
    """
     table の各セルの中に EL 形式があれば画像を挿入する
    """
    cell_top = shape.top
    for row_idx, row in enumerate(shape.table.rows):
        cell_left = shape.left
        for col_idx, col in enumerate(shape.table.columns):
            cell = row.cells[col_idx]

            if cell.is_merge_origin:
                # セル結合の場合
                height = 0
                for i in range(cell.span_height):
                    height += shape.table.rows[row_idx + i].height
                width = 0
                for i in range(cell.span_width):
                    width += shape.table.columns[col_idx + i].width
            else:
                height = row.height
                width = col.width

            rect = {'left': cell_left, 'top': cell_top, 'height': height, 'width': width}
            replace_img_in_text_frame(cell.text_frame, model, slide, rect)

            cell_left += col.width

        cell_top += row.height


def replace_img_in_shape(shape, model, slide):
    """
     shape の中のテキストに EL 形式があれば画像を挿入する
    """
    rect = {'left': shape.left, 'top': shape.top, 'height': shape.height, 'width': shape.width}
    replace_img_in_text_frame(shape.text_frame, model, slide, rect)


def replace_img_in_text_frame(text_frame, model, slide, rect):
    """
     text_frame 中のテキストに EL 形式があれば、それを model の該当する画像を挿入する
     見つからない場合は ValueError 例外となる
    """
    el = txt.extract_image(text_frame.text)
    if el is not None:
        value = pyel.eval_el(el, model)
        if not value:
            txt.replace_el_in_text_frame_with_str(text_frame, 'image:' + el, '')
            return
        elif isinstance(value, string_types):
            replace_img(slide, value, rect, Pt(5))
        elif isinstance(value, list):
            if len(value) == 0:
                txt.replace_el_in_text_frame_with_str(text_frame, 'image:' + el, '')
                return
            replace_imgs(slide, value, rect, Pt(5))
        else:
            raise ValueError(u"Invalid value for {%s}, model: %s" % (el, str(value)))

        text_frame.text = ''


def replace_imgs(slide, img_paths, rect, margin):
    """
     slide の指定された位置に複数の画像を挿入する
     画像が見つからない場合は ValueError 例外となる
    """
    count = len(img_paths)
    if count == 1:
        # 1枚の場合は全体に描画
        width = rect.get('width')
        height = rect.get('height')
    elif count == 2:
        # 2枚の場合は幅半分にそれぞれ描画
        width = rect.get('width') / 2
        height = rect.get('height')
    elif count >= 3:
        # 3枚以上の場合は縦横半分にそれぞれ描画
        width = rect.get('width') / 2
        height = rect.get('height') / 2
    else:
        raise ValueError(u"Invalid Pictures: %s" % str(img_paths))

    for i in range(count):
        img_rect = rect.copy()
        img_rect['height'] = height
        img_rect['width'] = width

        if i == 1:
            # 2枚目
            img_rect['left'] += width
        elif i == 2:
            # 3枚目
            img_rect['top'] += height
        elif i == 3:
            # 4枚目
            img_rect['left'] += width
            img_rect['top'] += height
        elif i == 4:
            # ToDo: とりあえず4枚まで
            break

        replace_img(slide, img_paths[i], img_rect, margin)


def replace_img(slide, img_path, rect, margin):
    """
     slide の指定された位置に画像を挿入する
     画像が見つからない場合は ValueError 例外となる
    """
    if not os.path.isfile(img_path):
        raise ValueError(u"Not found Picture: %s" % img_path)

    # スライドに画像を追加
    picture = slide.shapes.add_picture(img_path, 0, 0)
    if not isinstance(picture, Picture):
        raise ValueError(u"Invalid Picture: %s" % img_path)

    if float(picture.width) / rect.get('width') > float(picture.height) / rect.get('height'):
        width = rect.get('width') - (margin * 2)
        height = picture.height * width / picture.width
        left = rect.get('left') + margin
        top = rect.get('top') + (rect.get('height') - height) / 2
    else:
        height = rect.get('height') - (margin * 2)
        width = picture.width * height / picture.height
        left = rect.get('left') + (rect.get('width') - width) / 2
        top = rect.get('top') + margin

    # 画像を移動
    picture.left = int(left)
    picture.top = int(top)
    picture.width = int(width)
    picture.height = int(height)
