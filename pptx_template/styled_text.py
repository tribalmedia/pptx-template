#
# coding=utf-8

import logging
import re
import copy

import pptx_template.pyel as pyel
import pptx_template.pptx_util as util

from pptx.dml.color import RGBColor
from pptx.util import Pt

log = logging.getLogger()

_STYLED_RE = re.compile(r"\{styled_text:([A-Za-z0-9._\-]+)\}")


def _iterate_els(text):
    pos = 0
    while pos < len(text):
        text_id_match = _STYLED_RE.search(text[pos:])
        if text_id_match:
            pos = pos + text_id_match.end(1) + 1
            yield text_id_match.group(1)
        else:
            break


def _el_to_placeholder(el):
    return u"{styled_text:%s}" % el


def replace_all_els_in_table(table, model):
    """
     table の各セルの中に EL 形式があれば、それを model の該当する値と置き換える
    """
    for cell in [cell for row in table.rows for cell in row.cells]:
        replace_all_els_in_text_frame(cell.text_frame, model)


def replace_el_in_text_frame_with_list(text_frame, el, texts):
    """
     text_frame の各 paragraph.run 中のテキストに指定の EL 形式があれば、それを texts で置き換える
    """
    placeholder = _el_to_placeholder(el)
    for paragraph in text_frame.paragraphs:
        if placeholder not in paragraph.text:
            continue

        ((start_run, start_pos), (end_run, end_pos)) = _find_el_position([r.text for r in paragraph.runs], el)
        base_runs = copy.deepcopy(paragraph.runs)
        paragraph.clear()

        for (i, base_run) in enumerate(base_runs):
            if i == start_run and i == end_run:
                if base_run.text[0:start_pos]:
                    run = util.add_styled_run(paragraph, base_run)
                    run.text = base_run.text[0:start_pos]

                _insert_styled_run(paragraph, base_run, texts)

                if base_run.text[end_pos + 1:]:
                    run = util.add_styled_run(paragraph, base_run)
                    run.text = base_run.text[end_pos + 1:]

            elif i == start_run:
                if base_run.text[0:start_pos]:
                    run = util.add_styled_run(paragraph, base_run)
                    run.text = base_run.text[0:start_pos]

                _insert_styled_run(paragraph, base_run, texts)

            elif i == end_run:
                if base_run.text[end_pos + 1:]:
                    run = util.add_styled_run(paragraph, base_run)
                    run.text = base_run.text[end_pos + 1:]

            elif i < start_run or i > end_run:
                run = util.add_styled_run(paragraph, base_run)
                run.text = base_run.text

        return True
    return False


def replace_all_els_in_text_frame(text_frame, model):
    """
     text_frame 中のテキストに EL 形式が一つ以上あれば、それを model の該当する値と置き換える
    """
    for el in _iterate_els(text_frame.text):
        value = pyel.eval_el(el, model)
        if isinstance(value, list):
            replace_el_in_text_frame_with_list(text_frame, el, value)
        else:
            raise ValueError(u"Invalid value for {%s}, model: %s" % (el, str(value)))


def _find_el_position(texts, el):
    """
    text の配列中に分かれて記述されている EL の、先頭の '\{' の位置と、最後の '\}' の位置を返す。
    それぞれの位置は (text_index, position_in_text) の形で返される。
    """
    placeholder = _el_to_placeholder(el)
    full_text = ''.join(texts)

    start_pos = full_text.find(placeholder)
    if start_pos < 0:
        raise ValueError(u"texts %s doesn't contain EL:%s" % (texts, el))
    end_pos = start_pos + len(placeholder) - 1

    start_run_pos = start_run_index = end_run_pos = end_run_index = -1
    for (run_index, text) in enumerate(texts):
        length = len(text)

        if 0 <= start_pos < length:
            start_run_index = run_index
            start_run_pos = start_pos
        start_pos = start_pos - length

        if 0 <= end_pos < length:
            end_run_index = run_index
            end_run_pos = end_pos
            break
        end_pos = end_pos - length

    return ((start_run_index, start_run_pos), (end_run_index, end_run_pos))


def _insert_styled_run(paragraph, base_run, texts):
    for text in texts:
        if 'value' not in text:
            continue

        run = util.add_styled_run(paragraph, base_run)
        run.text = text['value']
        if 'font' in text:
            font = run.font
            if 'name' in text['font'] and isinstance(text['font']['name'], str):
                font.name = text['font']['name']
            if 'size' in text['font'] and isinstance(text['font']['size'], int):
                font.size = Pt(text['font']['size'])
            if 'bold' in text['font'] and isinstance(text['font']['bold'], bool):
                font.bold = text['font']['bold']
            if 'italic' in text['font'] and isinstance(text['font']['italic'], bool):
                font.italic = text['font']['italic']
            if 'underline' in text['font'] and isinstance(text['font']['underline'], bool):
                font.underline = text['font']['underline']
            if 'color' in text['font'] and isinstance(text['font']['color'], (list, str)):
                r = int(text['font']['color'][0], 16)
                g = int(text['font']['color'][1], 16)
                b = int(text['font']['color'][2], 16)
                font.color.rgb = RGBColor(r, g, b)
